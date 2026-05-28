# -*- coding: utf-8 -*-
"""
生成"数智心伴"项目简介PPT
运行：conda run -n aicaretaker python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── 颜色常量 ────────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x08, 0x10, 0x28)   # 深夜蓝（背景）
BG_CARD     = RGBColor(0x12, 0x1E, 0x42)   # 卡片底色
ACCENT_BLUE = RGBColor(0x3A, 0x8E, 0xFF)   # 亮蓝
ACCENT_PURP = RGBColor(0x9B, 0x5C, 0xFF)   # 紫色
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD4, 0xE8)
GOLD        = RGBColor(0xFF, 0xD7, 0x6E)
GREEN_OK    = RGBColor(0x3D, 0xE3, 0x8C)

# ─── 工具函数 ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # 完全空白
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=BG_DARK):
    """整页背景色"""
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None,
             line_color=None, line_width=Pt(0)):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=Pt(18), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    """添加文字框"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return txb

def add_para(tf, text, font_size=Pt(15), bold=False,
             color=LIGHT_GRAY, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    """向文字框追加段落"""
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return p

def title_bar(slide, title, subtitle=None):
    """顶部标题条"""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BG_CARD)
    # 左侧紫色竖线装饰
    add_rect(slide, 0.3, 0.18, 0.06, 0.74, fill_color=ACCENT_PURP)
    add_text_box(slide, title, 0.55, 0.15, 9, 0.75,
                 font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, 0.55, 0.72, 10, 0.35,
                     font_size=Pt(13), color=ACCENT_BLUE)

def gradient_deco(slide):
    """右上角渐变装饰圆（用两个半透明圆形叠加模拟）"""
    # 大圆
    s = slide.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x3A)
    s.line.fill.background()
    # 小圆
    s2 = slide.shapes.add_shape(9, Inches(11.5), Inches(-0.5), Inches(2.8), Inches(2.8))
    s2.fill.solid(); s2.fill.fore_color.rgb = RGBColor(0x2A, 0x10, 0x55)
    s2.line.fill.background()

# ─── 第1页：封面 ──────────────────────────────────────────────────────────────
def slide_cover(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)

    # 顶部亮蓝横线
    add_rect(sl, 0, 0, 13.33, 0.08, fill_color=ACCENT_BLUE)
    # 底部紫色横线
    add_rect(sl, 0, 7.42, 13.33, 0.08, fill_color=ACCENT_PURP)

    # 主标题
    add_text_box(sl, "数智心伴", 1.2, 1.5, 10, 1.6,
                 font_size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # 副标题
    add_text_box(sl, "多模态大模型驱动的AI数字人情感陪护系统",
                 1.2, 3.1, 11, 0.7,
                 font_size=Pt(22), bold=False, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)

    # 分割线
    add_rect(sl, 1.2, 3.9, 6, 0.04, fill_color=ACCENT_PURP)

    # 赛事信息
    add_text_box(sl, "中国大学生服务外包创新创业大赛  |  企业命题类（A类）",
                 1.2, 4.1, 10, 0.5,
                 font_size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    add_text_box(sl, "团队编号：___  |  团队名称：___  |  江苏大学",
                 1.2, 4.65, 10, 0.5,
                 font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    add_text_box(sl, "2026年4月", 1.2, 5.2, 4, 0.4,
                 font_size=Pt(13), color=RGBColor(0x88,0x99,0xBB), align=PP_ALIGN.LEFT)

    # 右侧装饰文字
    add_text_box(sl, "AI · 数字人 · 情感陪护",
                 8.5, 5.5, 4.5, 0.6,
                 font_size=Pt(14), color=ACCENT_PURP, align=PP_ALIGN.RIGHT)

# ─── 第2页：目录 ──────────────────────────────────────────────────────────────
def slide_toc(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "目  录", "CONTENTS")

    items = [
        ("01", "背景与痛点"),
        ("02", "项目定位与目标"),
        ("03", "系统总体架构"),
        ("04", "核心技术方案"),
        ("05", "功能演示与系统亮点"),
        ("06", "技术指标达成"),
        ("07", "应用场景与商业价值"),
        ("08", "团队介绍与分工"),
        ("09", "可行性分析"),
        ("10", "总结与展望"),
    ]

    # 两列布局
    for i, (num, text) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.4
        y = 1.35 + row * 1.05

        add_rect(sl, x, y, 5.9, 0.85, fill_color=BG_CARD,
                 line_color=ACCENT_BLUE, line_width=Pt(0.5))
        add_text_box(sl, num, x + 0.12, y + 0.1, 0.7, 0.65,
                     font_size=Pt(22), bold=True, color=ACCENT_PURP)
        add_text_box(sl, text, x + 0.8, y + 0.18, 5.0, 0.5,
                     font_size=Pt(17), bold=False, color=WHITE)

# ─── 第3页：背景与痛点 ────────────────────────────────────────────────────────
def slide_bg(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "背景与痛点", "BACKGROUND & PAIN POINTS")

    stats = [
        ("3亿+", "60岁以上老年人口"),
        ("50%+", "空巢独居老人占比"),
        ("98.3%", "失能老人有持续情感支持需求"),
    ]
    for i, (num, desc) in enumerate(stats):
        x = 0.5 + i * 4.25
        add_rect(sl, x, 1.25, 3.9, 1.55, fill_color=BG_CARD,
                 line_color=ACCENT_PURP, line_width=Pt(0.8))
        add_text_box(sl, num, x + 0.2, 1.35, 3.5, 0.85,
                     font_size=Pt(36), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, x + 0.1, 2.1, 3.7, 0.55,
                     font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    pain_points = [
        ("😔 情感陪伴严重缺失", "空巢、独居老人长期精神孤独，家庭陪护因代际分离日趋弱化"),
        ("🏥 专业资源极度稀缺", "心理咨询师严重供需失衡，专业疏导服务难以普及到普通家庭"),
        ("🤖 现有AI产品缺温度", "市场问答机器人缺乏情感深度，无法真正理解和回应情绪需求"),
    ]
    for i, (title_t, body) in enumerate(pain_points):
        y = 2.95 + i * 1.3
        add_rect(sl, 0.5, y, 12.3, 1.15, fill_color=BG_CARD)
        add_rect(sl, 0.5, y, 0.06, 1.15, fill_color=ACCENT_BLUE)
        add_text_box(sl, title_t, 0.75, y + 0.08, 4.5, 0.45,
                     font_size=Pt(15), bold=True, color=WHITE)
        add_text_box(sl, body, 0.75, y + 0.52, 11.8, 0.5,
                     font_size=Pt(12), color=LIGHT_GRAY)

# ─── 第4页：项目定位 ──────────────────────────────────────────────────────────
def slide_goal(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "项目定位与目标", "POSITIONING & OBJECTIVES")

    add_text_box(sl, "💡  愿景：让每一位老人都有一个懂自己的数字伙伴",
                 0.5, 1.2, 12.3, 0.65,
                 font_size=Pt(17), bold=True, color=GOLD)

    cores = [
        ("可陪伴", "7×24小时不间断\n有温度的情感交流", ACCENT_BLUE),
        ("可引导", "基于心理学知识库\n专业疏导焦虑抑郁情绪", ACCENT_PURP),
        ("可持续", "长期记忆·个性化\n越用越懂你", GREEN_OK),
    ]
    for i, (t, b, c) in enumerate(cores):
        x = 0.7 + i * 4.1
        add_rect(sl, x, 2.05, 3.7, 2.5, fill_color=BG_CARD, line_color=c, line_width=Pt(1.2))
        add_rect(sl, x, 2.05, 3.7, 0.12, fill_color=c)
        add_text_box(sl, t, x, 2.25, 3.7, 0.65,
                     font_size=Pt(24), bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text_box(sl, b, x + 0.15, 2.95, 3.4, 1.4,
                     font_size=Pt(14), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "核心闭环：感知  →  认知  →  干预  →  再感知",
                 1.5, 4.75, 10.3, 0.65,
                 font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.5, 4.72, 10.3, 0.72, fill_color=RGBColor(0x1A,0x0A,0x3A),
             line_color=ACCENT_PURP, line_width=Pt(0.8))
    # 覆盖修正层级——重新添加文字在矩形之上
    add_text_box(sl, "核心闭环：感知  →  认知  →  干预  →  再感知",
                 1.5, 4.78, 10.3, 0.65,
                 font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── 第5页：系统总体架构 ──────────────────────────────────────────────────────
def slide_arch(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "系统总体架构", "SYSTEM ARCHITECTURE")

    layers = [
        ("用户交互层",   "摄像头（视频）+ 麦克风（音频）+ 屏幕/扬声器（数字人呈现）",  ACCENT_BLUE),
        ("数据感知层",   "视觉：MediaPipe FaceMesh（478关键点·AU特征）  |  听觉：FunASR（Paraformer-zh · VAD）", RGBColor(0x3A,0xC8,0xFF)),
        ("认知决策层",   "多模态融合引擎  +  Qwen LLM对话  +  RAG心理知识库（ChromaDB）  +  10轮+上下文记忆", ACCENT_PURP),
        ("表达驱动层",   "CosyVoice TTS（流式）  +  口型/表情/动作驱动  +  Live2D / VRM 3D 实时渲染",          GOLD),
        ("数据层",       "SQLite用户档案  +  ChromaDB向量库  +  会话历史缓存  +  数字人资产库",                  GREEN_OK),
    ]

    for i, (name, desc, c) in enumerate(layers):
        y = 1.22 + i * 1.08
        add_rect(sl, 0.4, y, 0.18, 0.82, fill_color=c)
        add_rect(sl, 0.6, y, 12.3, 0.82, fill_color=BG_CARD)
        add_text_box(sl, name, 0.75, y + 0.1, 2.0, 0.6,
                     font_size=Pt(14), bold=True, color=c)
        add_text_box(sl, desc, 2.75, y + 0.15, 10.0, 0.55,
                     font_size=Pt(12), color=LIGHT_GRAY)
        if i < len(layers) - 1:
            add_text_box(sl, "↓", 6.5, y + 0.82, 0.5, 0.25,
                         font_size=Pt(10), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# ─── 第6页：核心技术——多模态感知 ─────────────────────────────────────────────
def slide_tech1(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "核心技术 · 多模态感知", "MULTIMODAL PERCEPTION")

    items = [
        ("👁  视觉感知", ACCENT_BLUE, [
            "MediaPipe FaceMesh：478个面部关键点实时检测",
            "提取15维AU（面部动作单元）情绪特征向量",
            "WebSocket 15fps 毫秒级实时传输后端",
            "自研FaceBehaviorModel：情绪→数字人驱动参数",
        ]),
        ("🎙  听觉感知", ACCENT_PURP, [
            "FunASR Paraformer-zh：中文端到端语音识别",
            "FSMN-VAD：精准语音活动检测",
            "CT-Punc：标点符号自动恢复",
            "WER ≤ 10%  |  SER ≤ 40%  达标赛题指标",
        ]),
    ]
    for i, (title_t, c, bullets) in enumerate(items):
        x = 0.5 + i * 6.4
        add_rect(sl, x, 1.25, 6.1, 5.8, fill_color=BG_CARD, line_color=c, line_width=Pt(0.8))
        add_rect(sl, x, 1.25, 6.1, 0.55, fill_color=c)
        add_text_box(sl, title_t, x + 0.2, 1.3, 5.7, 0.45,
                     font_size=Pt(18), bold=True, color=BG_DARK)
        for j, b in enumerate(bullets):
            add_text_box(sl, "▸  " + b, x + 0.25, 2.0 + j * 0.95, 5.6, 0.8,
                         font_size=Pt(14), color=LIGHT_GRAY)

# ─── 第7页：核心技术——LLM+RAG ─────────────────────────────────────────────────
def slide_tech2(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "核心技术 · LLM + RAG 情感对话引擎", "DIALOGUE ENGINE")

    flow = ["用户语音/文字输入", "RAG心理知识库检索\n（ChromaDB向量）",
            "Qwen大语言模型\n（8K+上下文）", "数字人回复生成\n+TTS语音输出"]
    colors = [ACCENT_BLUE, ACCENT_PURP, GOLD, GREEN_OK]
    for i, (txt, c) in enumerate(zip(flow, colors)):
        x = 0.5 + i * 3.1
        add_rect(sl, x, 1.3, 2.8, 1.35, fill_color=BG_CARD, line_color=c, line_width=Pt(1))
        add_text_box(sl, txt, x + 0.1, 1.45, 2.6, 1.05,
                     font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text_box(sl, "→", x + 2.82, 1.7, 0.35, 0.5,
                         font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    features = [
        ("🧠  Qwen大语言模型", "通义千问API，8K+上下文窗口，支持联网搜索增强，情感对话生成"),
        ("📚  RAG心理知识库", "ChromaDB向量检索，覆盖焦虑·抑郁·双向情感障碍等专业心理知识"),
        ("💾  长期上下文记忆", "10轮以上连续对话，会话摘要+记忆检索，个性化持续陪护"),
        ("🛡  专业知识兜底", "RAG无结果时动态启用LLM联网搜索，确保回复的专业性与安全性"),
    ]
    for i, (t, b) in enumerate(features):
        row, col = divmod(i, 2)
        x = 0.5 + col * 6.35
        y = 2.9 + row * 1.55
        add_rect(sl, x, y, 6.1, 1.35, fill_color=BG_CARD)
        add_rect(sl, x, y, 0.06, 1.35, fill_color=ACCENT_PURP)
        add_text_box(sl, t, x + 0.2, y + 0.1, 5.8, 0.45,
                     font_size=Pt(14), bold=True, color=WHITE)
        add_text_box(sl, b, x + 0.2, y + 0.55, 5.8, 0.65,
                     font_size=Pt(12), color=LIGHT_GRAY)

# ─── 第8页：数字人表达 ────────────────────────────────────────────────────────
def slide_tech3(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "核心技术 · 虚拟数字人表达与驱动", "DIGITAL HUMAN EXPRESSION")

    modes = [
        ("2D 数字人", ACCENT_BLUE, [
            "Live2D Cubism SDK + pixi-live2d-display",
            "PixiJS 图形渲染底层",
            "PARAM_* 参数实时驱动表情/动作",
            "端坐静止姿势·禁用自动Idle防抖动",
        ]),
        ("3D 数字人", ACCENT_PURP, [
            "Three.js r150 + @pixiv/three-vrm 2.x",
            "VRM 标准三维人形模型格式",
            "BlendShape 驱动面部表情",
            "GLTF 模型嘴部动画实时驱动",
        ]),
        ("TTS 语音合成", GOLD, [
            "CosyVoice 阿里开源神经TTS",
            "流式首包输出，降低感知延迟",
            "音素边界驱动口型同步（Lip-sync）",
            "Web Speech API 浏览器端备用方案",
        ]),
    ]
    for i, (title_t, c, bullets) in enumerate(modes):
        x = 0.4 + i * 4.3
        add_rect(sl, x, 1.25, 4.1, 5.75, fill_color=BG_CARD, line_color=c, line_width=Pt(0.8))
        add_rect(sl, x, 1.25, 4.1, 0.5, fill_color=c)
        add_text_box(sl, title_t, x + 0.15, 1.29, 3.8, 0.42,
                     font_size=Pt(16), bold=True, color=BG_DARK)
        for j, b in enumerate(bullets):
            add_text_box(sl, "▸  " + b, x + 0.2, 1.9 + j * 1.05, 3.7, 0.85,
                         font_size=Pt(13), color=LIGHT_GRAY)

# ─── 第9页：功能演示 ──────────────────────────────────────────────────────────
def slide_demo(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "功能演示与系统亮点", "FEATURE HIGHLIGHTS")

    cards = [
        ("💬 主对话界面", "数字人+历史对话侧边栏\n用户登录·会话持久化"),
        ("👁 视觉监控界面", "AU数据实时变化展示\nMediaPipe人脸追踪状态"),
        ("📚 心理知识库管理", "ChromaDB知识条目可视化\n支持实时查询与检索"),
        ("📊 WER/SER评测界面", "词错率/句错率联合评测\n达标赛题全部指标"),
    ]
    for i, (title_t, desc) in enumerate(cards):
        row, col = divmod(i, 2)
        x = 0.5 + col * 6.4
        y = 1.25 + row * 2.8
        add_rect(sl, x, y, 6.1, 2.55, fill_color=BG_CARD,
                 line_color=ACCENT_BLUE, line_width=Pt(0.6))
        # 图片占位框
        add_rect(sl, x + 0.1, y + 0.1, 5.9, 1.55,
                 fill_color=RGBColor(0x18, 0x28, 0x48))
        add_text_box(sl, "[ 截图占位 ]", x + 0.1, y + 0.55, 5.9, 0.6,
                     font_size=Pt(14), color=RGBColor(0x44,0x66,0x99),
                     align=PP_ALIGN.CENTER)
        add_text_box(sl, title_t, x + 0.15, y + 1.72, 5.8, 0.42,
                     font_size=Pt(14), bold=True, color=WHITE)
        add_text_box(sl, desc, x + 0.15, y + 2.1, 5.8, 0.4,
                     font_size=Pt(11), color=LIGHT_GRAY)

# ─── 第10页：技术指标 ─────────────────────────────────────────────────────────
def slide_metrics(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "技术指标达成情况", "TECHNICAL METRICS")

    metrics = [
        ("WER 词错率",   "≤ 10%",   "~8%",    "✅ 达标", ACCENT_BLUE),
        ("SER 句错率",   "≤ 40%",   "~15%",   "✅ 达标", ACCENT_PURP),
        ("连续对话轮次", "≥ 10轮",  "10轮+",  "✅ 达标", GREEN_OK),
        ("LLM上下文",    "≥ 8K tokens", "8K+",  "✅ 达标", GOLD),
        ("端到端响应",   "≤ 60秒",  "~30秒",  "✅ 达标", RGBColor(0xFF,0x88,0x44)),
    ]
    headers = ["指  标", "要  求", "实测结果", "状  态"]
    col_x = [0.5, 4.0, 7.0, 10.2]
    col_w = [3.4, 2.8, 3.0, 2.8]

    # 表头
    add_rect(sl, 0.5, 1.25, 12.3, 0.65, fill_color=RGBColor(0x1A,0x2E,0x5A))
    for j, h in enumerate(headers):
        add_text_box(sl, h, col_x[j], 1.3, col_w[j], 0.55,
                     font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (name, req, actual, status, c) in enumerate(metrics):
        y = 1.95 + i * 0.98
        bg = BG_CARD if i % 2 == 0 else RGBColor(0x0E, 0x18, 0x36)
        add_rect(sl, 0.5, y, 12.3, 0.88, fill_color=bg)
        add_rect(sl, 0.5, y, 0.08, 0.88, fill_color=c)
        vals = [name, req, actual, status]
        for j, v in enumerate(vals):
            clr = GREEN_OK if j == 3 else (WHITE if j == 0 else LIGHT_GRAY)
            add_text_box(sl, v, col_x[j], y + 0.18, col_w[j], 0.52,
                         font_size=Pt(15), bold=(j == 0),
                         color=clr, align=PP_ALIGN.CENTER)

# ─── 第11页：应用场景 ─────────────────────────────────────────────────────────
def slide_scenario(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "应用场景与商业价值", "SCENARIOS & BUSINESS VALUE")

    scenarios = [
        ("🏠 居家老人\n日常情感陪伴", "家庭智能终端接入\n7×24小时温暖陪伴\n主动感知情绪波动", ACCENT_BLUE),
        ("🏥 养老机构\n辅助护理", "替代部分人工陪护工作\n扩大服务覆盖面\n降低机构运营成本", ACCENT_PURP),
        ("⚠️ 心理健康\n预警监测", "异常情绪实时预警\n通知家属/医护人员\n焦虑抑郁早发现早干预", GOLD),
    ]
    for i, (title_t, desc, c) in enumerate(scenarios):
        x = 0.5 + i * 4.25
        add_rect(sl, x, 1.25, 3.9, 3.5, fill_color=BG_CARD, line_color=c, line_width=Pt(1))
        add_rect(sl, x, 1.25, 3.9, 0.12, fill_color=c)
        add_text_box(sl, title_t, x + 0.15, 1.4, 3.6, 1.0,
                     font_size=Pt(16), bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, x + 0.2, 2.45, 3.5, 2.0,
                     font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "市场规模：中国养老服务市场规模突破 10 万亿元",
                 0.5, 4.95, 8.5, 0.6,
                 font_size=Pt(16), bold=True, color=WHITE)
    add_text_box(sl, "盈利模式：B端（机构授权）+ C端（家庭订阅）+ G端（政府采购）",
                 0.5, 5.6, 12.3, 0.55,
                 font_size=Pt(14), color=LIGHT_GRAY)

# ─── 第12页：可行性分析 ───────────────────────────────────────────────────────
def slide_feasibility(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "可行性分析", "FEASIBILITY ANALYSIS")

    pros = [
        ("✅ 技术成熟可复现", "FunASR/Qwen/MediaPipe均为成熟方案，已完成原型系统联调"),
        ("✅ 低门槛部署",     "普通PC+摄像头即可运行，浏览器访问无需额外客户端"),
        ("✅ 指标全面达标",   "WER/SER/连续对话/响应时间5项赛题指标均已达标"),
    ]
    risks = [
        ("⚡ 隐私安全风险", "提供本地部署选项，音视频数据不上云"),
        ("⚡ 使用门槛风险", "极简UI设计，语音优先交互，适老化改造"),
        ("⚡ 模型幻觉风险", "RAG知识库兜底，专业心理内容审核过滤"),
    ]

    add_text_box(sl, "可行性亮点", 0.5, 1.2, 5.8, 0.55,
                 font_size=Pt(16), bold=True, color=ACCENT_BLUE)
    for i, (t, b) in enumerate(pros):
        y = 1.8 + i * 1.45
        add_rect(sl, 0.5, y, 5.9, 1.25, fill_color=BG_CARD)
        add_rect(sl, 0.5, y, 0.07, 1.25, fill_color=GREEN_OK)
        add_text_box(sl, t, 0.72, y + 0.1, 5.5, 0.45,
                     font_size=Pt(14), bold=True, color=WHITE)
        add_text_box(sl, b, 0.72, y + 0.58, 5.5, 0.55,
                     font_size=Pt(12), color=LIGHT_GRAY)

    add_text_box(sl, "风险应对", 7.0, 1.2, 5.8, 0.55,
                 font_size=Pt(16), bold=True, color=GOLD)
    for i, (t, b) in enumerate(risks):
        y = 1.8 + i * 1.45
        add_rect(sl, 7.0, y, 5.9, 1.25, fill_color=BG_CARD)
        add_rect(sl, 7.0, y, 0.07, 1.25, fill_color=GOLD)
        add_text_box(sl, t, 7.22, y + 0.1, 5.5, 0.45,
                     font_size=Pt(14), bold=True, color=WHITE)
        add_text_box(sl, b, 7.22, y + 0.58, 5.5, 0.55,
                     font_size=Pt(12), color=LIGHT_GRAY)

# ─── 第13页：团队介绍 ─────────────────────────────────────────────────────────
def slide_team(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "团队介绍与分工", "TEAM INTRODUCTION")

    roles = [
        ("🤖 AI算法/后端", "LLM对话引擎\nRAG知识库\n多模态融合"),
        ("🎙 语音方向",    "FunASR集成优化\nTTS语音合成\nWER/SER评测"),
        ("👁 视觉方向",    "MediaPipe感知\n面部驱动模型\nAU特征提取"),
        ("🎨 前端/数字人", "Vue3前端\nLive2D/VRM渲染\n口型动作同步"),
        ("📋 产品/文档",   "需求分析\n文档体系\nPPT设计"),
    ]
    for i, (role, tasks) in enumerate(roles):
        x = 0.3 + (i % 3) * 4.3 if i < 3 else 2.15 + (i - 3) * 4.3
        y = 1.25 if i < 3 else 3.95
        w = 3.9
        add_rect(sl, x, y, w, 2.45, fill_color=BG_CARD,
                 line_color=ACCENT_PURP, line_width=Pt(0.6))
        add_rect(sl, x, y, w, 0.5, fill_color=RGBColor(0x25,0x10,0x55))
        add_text_box(sl, role, x + 0.15, y + 0.05, w - 0.2, 0.4,
                     font_size=Pt(14), bold=True, color=ACCENT_PURP)
        add_text_box(sl, "姓名：___（占位）", x + 0.15, y + 0.6, w - 0.2, 0.4,
                     font_size=Pt(12), color=GOLD)
        add_text_box(sl, tasks, x + 0.15, y + 1.05, w - 0.2, 1.2,
                     font_size=Pt(12), color=LIGHT_GRAY)

    add_text_box(sl, "指导教师：_________（占位）  |  学校：江苏大学",
                 0.5, 7.05, 12.3, 0.38,
                 font_size=Pt(12), color=RGBColor(0x66,0x88,0xBB))

# ─── 第14页：总结与展望 ───────────────────────────────────────────────────────
def slide_summary(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gradient_deco(sl)
    title_bar(sl, "总结与展望", "SUMMARY & OUTLOOK")

    add_text_box(sl,
        "数智心伴  ——  构建有温度的 AI 数字人情感陪护闭环系统",
        0.5, 1.25, 12.3, 0.75,
        font_size=Pt(20), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    achieved = [
        "多模态感知：MediaPipe视觉 + FunASR听觉实时采集",
        "LLM对话引擎：Qwen + RAG心理知识库专业对话",
        "数字人渲染：2D Live2D / 3D VRM双模高表现力",
        "全指标达标：WER/SER/响应时间/连续对话均通过",
    ]
    add_text_box(sl, "✅  已实现（当前版本）", 0.5, 2.15, 5.8, 0.5,
                 font_size=Pt(15), bold=True, color=GREEN_OK)
    for i, a in enumerate(achieved):
        add_text_box(sl, "▸  " + a, 0.5, 2.75 + i * 0.82, 5.8, 0.65,
                     font_size=Pt(13), color=LIGHT_GRAY)

    future = [
        "接入PHQ-9/GAD-7等心理评估量表",
        "支持移动端（小程序/App）",
        "融合智能穿戴设备（心率/睡眠）",
        "探索 AI + 人工 混合陪护模式",
    ]
    add_text_box(sl, "🚀  未来规划", 7.0, 2.15, 5.8, 0.5,
                 font_size=Pt(15), bold=True, color=ACCENT_BLUE)
    for i, f in enumerate(future):
        add_text_box(sl, "▸  " + f, 7.0, 2.75 + i * 0.82, 5.8, 0.65,
                     font_size=Pt(13), color=LIGHT_GRAY)

    add_rect(sl, 0.5, 6.5, 12.3, 0.75, fill_color=BG_CARD)
    add_text_box(sl,
        "感谢评审老师的宝贵时间  |  数智心伴团队  敬上",
        0.5, 6.55, 12.3, 0.6,
        font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── 主程序 ───────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()

    slide_cover(prs)
    slide_toc(prs)
    slide_bg(prs)
    slide_goal(prs)
    slide_arch(prs)
    slide_tech1(prs)
    slide_tech2(prs)
    slide_tech3(prs)
    slide_demo(prs)
    slide_metrics(prs)
    slide_scenario(prs)
    slide_feasibility(prs)
    slide_team(prs)
    slide_summary(prs)

    out = r"d:\AI数字人情感陪护项目\ppt制作指导说明\数智心伴_项目简介PPT.pptx"
    prs.save(out)
    print(f"PPT已生成：{out}")
    print(f"共 {len(prs.slides)} 页")

if __name__ == "__main__":
    main()
